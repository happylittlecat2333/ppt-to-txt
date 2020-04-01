from pptx import Presentation
import glob

for eachfile in glob.glob("*.pptx"):
    
    print("正在转换:",eachfile)
    prs = Presentation(eachfile)

    f = open(eachfile[:-5]+ ".txt","w")
    for slide in prs.slides: 
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
#                     text_runs.append(run.text)
                    f.write("{}".format(run.text))
                f.write("\n")

    f.close()
    print("成功转换:",eachfile[:-5]+ ".txt")